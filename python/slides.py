from pptx import Presentation
import argparse
import csv
import sys
import csv
import os



originalDir=sys.argv[1]
if not os.path.exists(originalDir+"/Presentation/presentation.pptx"):
    sys.exit()
prs = Presentation(originalDir+"/Presentation/presentation.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []
slideErrors = [None] * len(prs.slides)
slides=0
for slide in prs.slides:
    error=""
    slides=slides+1
    shapes=0
    textSlide=""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        shapes=shapes+1
        textSlide=textSlide+shape.text
        paragraphs=0
        for paragraph in shape.text_frame.paragraphs:
            paragraphs=paragraphs+1
            #print("Font Size: ", paragraph.font.size)
            runs=0
            for run in paragraph.runs:
                runs=runs+1
                #print("Run ", runs)
                size=run.font.size
                if size is not None:
                    if size<22:
                        error="Small"
                    if size<18:
                        error="Tiny"
    if len(textSlide.split()) > 60:
        error=error+"Long"
    elif len(textSlide.split())  > 33:
        error=error+"Verbose"
    slideErrors[slides-1]=error
print(slideErrors)
resultFile = open(originalDir+"/Presentation/result.csv",'w')
wr = csv.writer(resultFile)
index=0

for errors in slideErrors:
    index=index+1
    if errors=="":
        errors="None"
    wr.writerow([index,errors])
resultFile.close()
